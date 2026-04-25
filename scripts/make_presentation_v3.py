"""
make_presentation_v3.py
Generate a 21-slide presentation with KAUST-style visual template:
  - Light-green top banner with title
  - Logo placeholder (top-right)
  - Bottom green strip + page number
  - Green geometric section dividers
Each content slide carries:
  - Visible bullets summarising what the speaker will say (so the deck reads
    standalone at a glance), plus the full English script in the notes panel.

Run with:
    D:/anaconda3/envs/hydroclimate/python.exe scripts/make_presentation_v3.py
Output:
    report/Presentation_v3.pptx
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG_DIR = os.path.join(ROOT, 'report', 'figures_fp')
OUT_PPTX = os.path.join(ROOT, 'report', 'Presentation_v3.pptx')

# ---------- palette ----------
BANNER_GREEN = RGBColor(0xD4, 0xF0, 0xDC)
DIVIDER_GREEN = RGBColor(0xC4, 0xEA, 0xCF)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
SAUDI = RGBColor(0xE0, 0x7B, 0x39)
ITALY = RGBColor(0x4A, 0x9D, 0x5C)
BANGLA = RGBColor(0x3F, 0x7C, 0xAC)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


# =========================================================================
# Low-level helpers
# =========================================================================
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, text, left, top, width, height,
            size=20, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = [text] if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def filled_rect(slide, left, top, width, height, color, line=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if not line:
        sh.line.fill.background()
    return sh


def filled_triangle(slide, left, top, width, height, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_image(slide, name, left, top, width=None, height=None):
    path = os.path.join(FIG_DIR, name)
    if not os.path.exists(path):
        textbox(slide, f'[Image placeholder: {name}]',
                left, top, Inches(5), Inches(0.5),
                size=12, color=GREY)
        return None
    kwargs = {}
    if width is not None:
        kwargs['width'] = width
    if height is not None:
        kwargs['height'] = height
    return slide.shapes.add_picture(path, left, top, **kwargs)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# =========================================================================
# Template chrome
# =========================================================================
def apply_chrome(slide, title_text, page_num):
    """Apply the standard top banner, logo placeholder, bottom strip,
    and page number to a content slide."""
    # Top banner
    filled_rect(slide, 0, 0, SLIDE_W, Inches(0.85), BANNER_GREEN)
    # Title text on banner
    textbox(slide, title_text, Inches(0.4), Inches(0.12),
            Inches(11.4), Inches(0.65),
            size=28, bold=True, color=DARK, font='Cambria')
    # Logo placeholder top right
    logo = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(12.55), Inches(0.18),
        Inches(0.5), Inches(0.5),
    )
    logo.fill.solid()
    logo.fill.fore_color.rgb = RGBColor(0xF5, 0xC0, 0x4B)
    logo.line.color.rgb = RGBColor(0x9E, 0x6E, 0x10)
    # Bottom strip
    filled_rect(slide, 0, SLIDE_H - Inches(0.2), SLIDE_W,
                Inches(0.2), BANNER_GREEN)
    # Page number
    textbox(slide, str(page_num), Inches(12.7), SLIDE_H - Inches(0.45),
            Inches(0.5), Inches(0.25),
            size=11, color=GREY, align=PP_ALIGN.RIGHT)


def section_divider_chrome(slide, label):
    """Section-divider slide layout: small green strip top, big green right
    triangle, large left-aligned title."""
    # Top thin strip
    filled_rect(slide, 0, 0, SLIDE_W, Inches(0.6), BANNER_GREEN)
    # Bottom thin strip
    filled_rect(slide, 0, SLIDE_H - Inches(0.4), SLIDE_W,
                Inches(0.4), BANNER_GREEN)
    # Right side triangle (geometric accent)
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, Inches(7.5), Inches(0.6),
        Inches(5.83), Inches(6.5),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = DIVIDER_GREEN
    sh.line.fill.background()
    sh.rotation = 90
    # Section title
    textbox(slide, label, Inches(0.6), Inches(2.8),
            Inches(7.5), Inches(1.5),
            size=58, bold=True, color=DARK, font='Cambria')


# =========================================================================
# Slide builders
# =========================================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---------- 1. Title ----------
    s = add_blank(prs)
    # Diagonal green panel left side
    sh = slide_diagonal = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, Inches(-0.5), Inches(-0.5),
        Inches(6), Inches(8),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = DIVIDER_GREEN
    sh.line.fill.background()
    # Top-right logo placeholder
    logo = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(12.4), Inches(0.4),
        Inches(0.7), Inches(0.7),
    )
    logo.fill.solid()
    logo.fill.fore_color.rgb = RGBColor(0xF5, 0xC0, 0x4B)
    logo.line.color.rgb = RGBColor(0x9E, 0x6E, 0x10)
    textbox(s, 'KAUST', Inches(11.0), Inches(0.55),
            Inches(1.3), Inches(0.4),
            size=13, bold=True, color=DARK, align=PP_ALIGN.RIGHT)
    # Title
    textbox(s, 'Climate Controls on Runoff Across Hydroclimatic Regimes',
            Inches(1), Inches(2.8), Inches(11.3), Inches(1.4),
            size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER,
            font='Cambria')
    textbox(s, 'Using ERA5-Land and Machine Learning Analysis',
            Inches(1), Inches(4.0), Inches(11.3), Inches(0.7),
            size=22, color=GREY, align=PP_ALIGN.CENTER, font='Cambria')
    textbox(s,
            ['Geo-Environmental Modeling and Analysis',
             'Mariia Solodiankina  ·  Shengli Zhu',
             'April 29, 2026'],
            Inches(1), Inches(5.4), Inches(11.3), Inches(1.6),
            size=18, color=DARK, align=PP_ALIGN.CENTER)
    set_notes(s,
              "Good morning, everyone. Our project is 'Climate Controls on Runoff Across "
              "Hydroclimatic Regimes Using ERA5-Land and Machine Learning'. I'm Mariia, "
              "and Shengli will join later for the results.")

    # ---------- 2. Outline ----------
    s = add_blank(prs)
    apply_chrome(s, 'Outline', 2)
    items = [
        ('01', 'Introduction'),
        ('02', 'Study Area'),
        ('03', 'Methodology'),
        ('04', 'Results'),
        ('05', 'Conclusion'),
    ]
    for i, (num, name) in enumerate(items):
        col = i % 3
        row = i // 3
        x = Inches(1.0 + col * 4.0)
        y = Inches(1.8 + row * 2.5)
        textbox(s, num, x, y, Inches(2), Inches(0.9),
                size=44, bold=True, color=DARK, font='Cambria')
        textbox(s, name, x, y + Inches(1.0), Inches(3), Inches(0.6),
                size=22, color=DARK, font='Cambria')
    set_notes(s,
              "Our talk has five parts: introduction, study area, methodology, "
              "results, and conclusion.")

    # ---------- 3. Intro divider ----------
    s = add_blank(prs)
    section_divider_chrome(s, 'Introduction')
    set_notes(s, "Let's start with the introduction.")

    # ---------- 4. Background ----------
    s = add_blank(prs)
    apply_chrome(s, 'Background', 4)
    textbox(s,
            ['•  Runoff is a key component of the water cycle',
             '',
             '•  Generation mechanism differs fundamentally across climate regimes',
             '     —  arid zones: episodic, event-driven pulses',
             '     —  humid zones: near-linear precipitation tracking',
             '',
             '•  Understanding these differences is essential for hydrological',
             '    prediction and water resource management'],
            Inches(0.6), Inches(1.3), Inches(7.2), Inches(5),
            size=20, color=DARK)
    textbox(s, '[ Place USGS water cycle figure here ]',
            Inches(8.0), Inches(2.6), Inches(4.8), Inches(2.5),
            size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    filled_rect(s, Inches(8.0), Inches(2.6), Inches(4.8), Inches(2.5),
                RGBColor(0xF8, 0xF8, 0xF8))
    textbox(s, '[ Place USGS water cycle figure here ]',
            Inches(8.0), Inches(3.6), Inches(4.8), Inches(0.5),
            size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    set_notes(s,
              "Runoff is a key component of the water cycle, but its generation mechanism "
              "is not the same across climates. In arid zones, runoff is driven by episodic "
              "pulses. In humid zones, it tracks precipitation almost linearly. "
              "Understanding these differences matters for hydrological prediction and "
              "water resource management — which motivates our study.")

    # ---------- 5. Introduction (Problem + RQs) ----------
    s = add_blank(prs)
    apply_chrome(s, 'Introduction', 5)
    textbox(s, 'Problem',
            Inches(0.6), Inches(1.3), Inches(2), Inches(0.6),
            size=22, bold=True, color=DARK)
    textbox(s,
            'The relative importance of climate drivers on runoff has not been '
            'systematically quantified across the aridity gradient.',
            Inches(2.7), Inches(1.3), Inches(10), Inches(1.2),
            size=18, color=DARK)
    textbox(s, 'Research Questions',
            Inches(0.6), Inches(3.0), Inches(3), Inches(0.6),
            size=22, bold=True, color=DARK)
    rqs = [
        '1.  How do water balance and runoff efficiency differ across hydroclimatic regimes?',
        '2.  What long-term hydroclimatic trends exist in these regions?',
        '3.  Which climate variables dominate runoff generation, and by how much does this dominance shift across regimes?',
    ]
    for i, rq in enumerate(rqs):
        textbox(s, rq, Inches(0.6), Inches(3.7 + i * 0.95),
                Inches(12), Inches(0.9),
                size=17, color=DARK)
    set_notes(s,
              "The core problem is this: although we know climate drivers of runoff differ "
              "between dry and wet regions, no one has systematically quantified how their "
              "relative importance changes across the aridity gradient. We address this with "
              "three research questions.")

    # ---------- 6. Study Area divider ----------
    s = add_blank(prs)
    section_divider_chrome(s, 'Study Area')
    set_notes(s, "Now let's look at the study area.")

    # ---------- 7. Study Area ----------
    s = add_blank(prs)
    apply_chrome(s, 'Study Area', 7)
    add_image(s, 'fig01_study_area.png',
              Inches(0.4), Inches(1.1), width=Inches(12.5))
    # Table-like data block at bottom
    headers = ['Country', 'Climate (Köppen)', 'P (mm yr⁻¹)', 'RC (R/P)']
    rows = [
        ('Saudi Arabia', 'hyper-arid (BWh)', '67', '0.03', SAUDI),
        ('Italy', 'transitional (Csa/Cfa)', '1,047', '0.41', ITALY),
        ('Bangladesh', 'tropical monsoon (Am)', '2,316', '0.56', BANGLA),
    ]
    cell_w = [Inches(2.6), Inches(3.5), Inches(2.6), Inches(2.6)]
    cell_x = [Inches(0.6)]
    for w in cell_w[:-1]:
        cell_x.append(cell_x[-1] + w)
    # Header row
    y0 = Inches(5.4)
    filled_rect(s, Inches(0.6), y0, sum(cell_w, Inches(0)),
                Inches(0.45), RGBColor(0x4A, 0x77, 0x4F))
    for h, x, w in zip(headers, cell_x, cell_w):
        textbox(s, h, x, y0, w, Inches(0.45),
                size=14, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    # Data rows
    for i, row in enumerate(rows):
        y = y0 + Inches(0.5 + i * 0.45)
        accent_col = row[4]
        # Country cell with light tint
        tint = RGBColor(
            min(255, accent_col[0] + 50),
            min(255, accent_col[1] + 50),
            min(255, accent_col[2] + 50),
        )
        filled_rect(s, cell_x[0], y, cell_w[0], Inches(0.4), tint)
        for j, (val, x, w) in enumerate(zip(row[:4], cell_x, cell_w)):
            textbox(s, val, x, y + Inches(0.05), w, Inches(0.35),
                    size=13, color=DARK, bold=(j == 0),
                    align=PP_ALIGN.CENTER)
    set_notes(s,
              "We selected three countries that span the global aridity gradient. "
              "Saudi Arabia is hyper-arid, with 67 millimetres per year of precipitation "
              "and a runoff coefficient of just three percent. Italy is the transitional "
              "case at 1047 millimetres, with 41 percent of precipitation becoming runoff. "
              "Bangladesh is humid tropical monsoon, 2316 millimetres, 56 percent runoff "
              "coefficient. Just from water balance alone, the runoff coefficient already "
              "varies by a factor of 20.")

    # ---------- 8. Methodology divider ----------
    s = add_blank(prs)
    section_divider_chrome(s, 'Methodology')
    set_notes(s, "Next, the methodology.")

    # ---------- 9. Methodology ----------
    s = add_blank(prs)
    apply_chrome(s, 'Methodology', 9)
    steps = [
        ('1', 'Data Collection',
         'ERA5-Land monthly  ·  1950–2025  ·  0.1° resolution'),
        ('2', 'Water Balance Analysis',
         'P = ET + R + ΔS + ε'),
        ('3', 'Trend Analysis',
         "Mann–Kendall test  +  Sen's slope  ·  annual / seasonal / pixel-level"),
        ('4', 'Machine Learning Runoff Prediction',
         'XGBoost regression  ·  8 climate features  ·  '
         'train 1950–2004 / test 2005–2025  ·  R², NSE, KGE'),
        ('5', 'SHAP Attribution',
         'TreeExplainer  ·  per-sample, model-exact decomposition'),
    ]
    for i, (num, head, body) in enumerate(steps):
        y = Inches(1.3 + i * 1.0)
        # Number circle
        circ = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.7), Inches(0.7),
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = DIVIDER_GREEN
        circ.line.fill.background()
        textbox(s, num, Inches(0.6), y + Inches(0.13),
                Inches(0.7), Inches(0.5),
                size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        # Head + body
        textbox(s, head, Inches(1.55), y, Inches(11), Inches(0.45),
                size=18, bold=True, color=DARK)
        textbox(s, body, Inches(1.55), y + Inches(0.45),
                Inches(11), Inches(0.55),
                size=14, color=GREY, font='Consolas')
    set_notes(s,
              "We follow a five-step pipeline. First, we collect ERA5-Land monthly data "
              "for 1950 to 2025 at 0.1 degree resolution. Second, we close the water balance "
              "with the equation P equals ET plus R plus delta-S plus a residual. Third, "
              "we test for long-term trends using Mann-Kendall and Sen's slope. Fourth, we "
              "train an XGBoost regression model using eight climate features as inputs and "
              "runoff as the target — split chronologically with 1950 to 2004 for training "
              "and 2005 to 2025 for testing — evaluated with R-squared, NSE, and KGE. Fifth, "
              "we apply SHAP to attribute the trained model's predictions to individual "
              "climate drivers.")

    # ---------- 10. Results divider ----------
    s = add_blank(prs)
    section_divider_chrome(s, 'Results')
    set_notes(s, "Now Shengli will walk us through the results.")

    # ---------- 11. Water balance — table + spatial ----------
    s = add_blank(prs)
    apply_chrome(s, 'Water balance — partitioning across climates', 11)
    # Table
    headers = ['Country', 'P', 'ET', 'R', 'RC']
    rows = [
        ('Saudi Arabia', '67', '72*', '2.5', '0.03', SAUDI),
        ('Italy', '1,047', '633', '435', '0.41', ITALY),
        ('Bangladesh', '2,316', '1,052', '1,317', '0.56', BANGLA),
    ]
    cell_w = [Inches(1.9), Inches(0.8), Inches(0.8), Inches(0.8), Inches(0.8)]
    cell_x = [Inches(0.5)]
    for w in cell_w[:-1]:
        cell_x.append(cell_x[-1] + w)
    y0 = Inches(1.3)
    filled_rect(s, Inches(0.5), y0, sum(cell_w, Inches(0)),
                Inches(0.45), RGBColor(0x4A, 0x77, 0x4F))
    for h, x, w in zip(headers, cell_x, cell_w):
        textbox(s, h, x, y0, w, Inches(0.45),
                size=14, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        y = y0 + Inches(0.5 + i * 0.45)
        tint = RGBColor(
            min(255, row[5][0] + 50),
            min(255, row[5][1] + 50),
            min(255, row[5][2] + 50),
        )
        filled_rect(s, cell_x[0], y, cell_w[0], Inches(0.4), tint)
        for j, (val, x, w) in enumerate(zip(row[:5], cell_x, cell_w)):
            textbox(s, val, x, y + Inches(0.05), w, Inches(0.35),
                    size=13, color=DARK, bold=(j == 0), align=PP_ALIGN.CENTER)
    textbox(s, 'Units: mm yr⁻¹     * ERA5-Land bias in hyper-arid zones (data artefact)',
            Inches(0.5), Inches(3.4), Inches(5.5), Inches(0.4),
            size=10, color=GREY)
    # Bullets
    textbox(s,
            ['•  Runoff coefficient: 3%  →  41%  →  56%',
             '•  ≈ 20× range in water-balance partitioning alone',
             '•  Spatial maps (right) reveal strong within-country heterogeneity:',
             '     —  Saudi P/R concentrate along the Hijaz mountains',
             '     —  Italy runoff peaks in the Alps',
             '     —  Bangladesh peaks in the Meghna flood zone (north-east)'],
            Inches(0.5), Inches(4.0), Inches(6.5), Inches(3),
            size=14, color=DARK)
    add_image(s, 'fig02_spatial_distribution.png',
              Inches(7.2), Inches(1.3), width=Inches(5.9))
    set_notes(s,
              "The water balance shows regime-defining differences. Saudi Arabia consumes "
              "nearly all 67 millimetres of annual rainfall through evapotranspiration; "
              "runoff coefficient is just three percent. Italy partitions roughly 41 percent "
              "of rainfall into runoff. Bangladesh exceeds 56 percent. The spatial maps on "
              "the right reveal within-country heterogeneity. Note that ET exceeds P in "
              "Saudi Arabia at 1.11 — this is a known ERA5-Land bias in hyper-arid regions "
              "and we treat it as a data artefact.")

    # ---------- 12. Water balance — timeseries ----------
    s = add_blank(prs)
    apply_chrome(s, 'Water balance — annual time series (1950–2025)', 12)
    add_image(s, 'fig03_water_balance_timeseries.png',
              Inches(0.4), Inches(1.1), width=Inches(8.8))
    # Right-side bullets with country colour dots
    bullets = [
        ('Saudi Arabia', SAUDI,
         'high inter-annual variability — event-driven runoff'),
        ('Italy', ITALY,
         'relatively stable, modest long-term shifts'),
        ('Bangladesh', BANGLA,
         'clearly declining P and R since ~2000'),
    ]
    for i, (name, col, body) in enumerate(bullets):
        y = Inches(1.5 + i * 1.8)
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(9.4), y + Inches(0.12),
            Inches(0.25), Inches(0.25),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = col
        dot.line.fill.background()
        textbox(s, name, Inches(9.75), y, Inches(3.4), Inches(0.5),
                size=16, bold=True, color=col)
        textbox(s, body, Inches(9.75), y + Inches(0.55),
                Inches(3.4), Inches(1.2),
                size=13, color=DARK)
    set_notes(s,
              "Looking at 76 years of annual time series, Saudi Arabia shows extremely high "
              "inter-annual variability — its runoff is driven by episodic events. Italy is "
              "relatively stable with modest long-term shifts. Bangladesh shows a clear "
              "declining trend in both precipitation and runoff since around 2000. This "
              "decline becomes statistically significant in the next slide.")

    # ---------- 13. Mann-Kendall trends ----------
    s = add_blank(prs)
    apply_chrome(s, 'Mann–Kendall trends (1950–2025, annual)', 13)
    # Table
    headers = ['Country', 'P', 'ET', 'R', 'T']
    rows = [
        ('Saudi Arabia', '—', '—', '↓***', '↑***', SAUDI),
        ('Italy', '—', '↑***', '—', '↑***', ITALY),
        ('Bangladesh', '↓**', '—', '↓**', '↑***', BANGLA),
    ]
    cell_w = [Inches(1.9), Inches(0.8), Inches(0.8), Inches(0.8), Inches(0.8)]
    cell_x = [Inches(0.5)]
    for w in cell_w[:-1]:
        cell_x.append(cell_x[-1] + w)
    y0 = Inches(1.3)
    filled_rect(s, Inches(0.5), y0, sum(cell_w, Inches(0)),
                Inches(0.45), RGBColor(0x4A, 0x77, 0x4F))
    for h, x, w in zip(headers, cell_x, cell_w):
        textbox(s, h, x, y0, w, Inches(0.45),
                size=14, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        y = y0 + Inches(0.5 + i * 0.45)
        tint = RGBColor(
            min(255, row[5][0] + 50),
            min(255, row[5][1] + 50),
            min(255, row[5][2] + 50),
        )
        filled_rect(s, cell_x[0], y, cell_w[0], Inches(0.4), tint)
        for j, (val, x, w) in enumerate(zip(row[:5], cell_x, cell_w)):
            textbox(s, val, x, y + Inches(0.05), w, Inches(0.35),
                    size=14, color=DARK, bold=(j == 0), align=PP_ALIGN.CENTER)
    textbox(s, '↑/↓ : sign of Sen\'s slope    ** p<0.05    *** p<0.01',
            Inches(0.5), Inches(3.4), Inches(6), Inches(0.4),
            size=11, color=GREY)
    # Bullets (sorted by message strength)
    bullet_blocks = [
        (BANGLA, 'Bangladesh:',
         'P & R declining together  →  weakening summer-monsoon signature'),
        (ITALY, 'Italy:',
         "ET +0.88 mm yr⁻¹ without matching P decline  →  warming fingerprint"),
        (SAUDI, 'Saudi Arabia:',
         'no peninsula-wide trend  →  signal confined to Hijaz–Asir highlands'),
    ]
    for i, (col, head, body) in enumerate(bullet_blocks):
        y = Inches(4.0 + i * 0.9)
        textbox(s, head, Inches(0.5), y, Inches(2.0), Inches(0.4),
                size=15, bold=True, color=col)
        textbox(s, body, Inches(2.5), y, Inches(4.5), Inches(0.7),
                size=13, color=DARK)
    add_image(s, 'fig08_spatial_trend_maps.png',
              Inches(7.1), Inches(1.3), width=Inches(6.0))
    set_notes(s,
              "All three countries warm significantly. Bangladesh shows simultaneous decline "
              "in both precipitation and runoff — a signature of weakening summer monsoon. "
              "Italy's evapotranspiration increases by 0.88 millimetres per year without a "
              "matching precipitation decline — this is the classic warming fingerprint in "
              "the water balance. Saudi Arabia shows no significant trend across most of the "
              "peninsula; the signal is confined to the Hijaz-Asir highlands. The pixel-level "
              "maps on the right reveal opposing trends within Italy.")

    # ---------- 14. Hydroclimatic seasonal trends ----------
    s = add_blank(prs)
    apply_chrome(s, 'Seasonal breakdown of trends', 14)
    add_image(s, 'fig07_trend_heatmap.png',
              Inches(0.3), Inches(1.1), width=Inches(12.7))
    bullets = [
        (ITALY, 'Italy ET',
         'increase strongest in spring (MAM, +1.99) and summer (JJA, +1.37)'),
        (BANGLA, 'Bangladesh P',
         'decline concentrated in MAM (−6.31) and JJA (−6.33) — core monsoon seasons'),
        (SAUDI, 'Saudi warming',
         'broad across seasons; autumn strongest (+0.05 °C yr⁻¹)'),
    ]
    for i, (col, head, body) in enumerate(bullets):
        y = Inches(5.5 + i * 0.5)
        textbox(s, head, Inches(0.4), y, Inches(2.4), Inches(0.4),
                size=14, bold=True, color=col)
        textbox(s, body, Inches(2.9), y, Inches(10), Inches(0.4),
                size=13, color=DARK)
    set_notes(s,
              "Breaking the trends down by season exposes details the annual table hides. "
              "Italy's ET increase is concentrated in spring and summer, when atmospheric "
              "demand is highest. Bangladesh's precipitation decline is sharpest in spring "
              "and summer — exactly the core monsoon seasons. Saudi warming is broad across "
              "all seasons but strongest in autumn. These seasonal asymmetries point to "
              "specific physical mechanisms at work.")

    # ---------- 15. Model Performance ----------
    s = add_blank(prs)
    apply_chrome(s, 'Model performance — predictability scales with humidity', 15)
    add_image(s, 'fig09b_annual_predicted_vs_observed.png',
              Inches(0.5), Inches(1.5), width=Inches(12.3))
    # KGE highlight box
    filled_rect(s, Inches(1.5), Inches(6.0), Inches(10.3),
                Inches(0.85), DIVIDER_GREEN)
    textbox(s,
            'Annual KGE rises from 0.69 (arid) → 0.96 (transitional) → 0.98 (humid)',
            Inches(1.5), Inches(6.05), Inches(10.3), Inches(0.4),
            size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    textbox(s,
            '— a direct reflection of the Budyko transition from water- to energy-limited regimes',
            Inches(1.5), Inches(6.45), Inches(10.3), Inches(0.4),
            size=13, color=GREY, align=PP_ALIGN.CENTER)
    set_notes(s,
              "Our XGBoost model's predictive skill increases monotonically with humidity. "
              "KGE rises from 0.69 in the arid case to 0.98 in the humid case. The physical "
              "interpretation: in arid regions, runoff is driven by stochastic extreme events "
              "that monthly or annual climate means cannot capture. In humid regions, runoff "
              "tracks precipitation almost linearly. This monotonic skill gradient is itself "
              "a reflection of the Budyko transition.")

    # ---------- 16. SHAP — 250x bar (KEY SLIDE) ----------
    s = add_blank(prs)
    apply_chrome(s, 'Dominant runoff drivers (SHAP)', 16)
    textbox(s, 'Precipitation dominates — by wildly different amounts',
            Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
            size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    textbox(s,
            'Each mm of annual P adds ~0 mm of runoff in Saudi Arabia, '
            '~1 mm in Bangladesh — non-linear amplification by saturation effects.',
            Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.5),
            size=13, color=GREY, align=PP_ALIGN.CENTER)
    # Bar chart placeholder area
    filled_rect(s, Inches(2.0), Inches(2.4), Inches(9.3),
                Inches(2.6), RGBColor(0xF8, 0xFB, 0xF8))
    textbox(s, '[ Place your custom 250× horizontal bar chart here ]',
            Inches(2.0), Inches(2.7), Inches(9.3), Inches(0.4),
            size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    # Three big numbers (visual approximation of the bars)
    textbox(s, 'Saudi Arabia',
            Inches(2.5), Inches(3.0), Inches(2.5), Inches(0.5),
            size=14, bold=True, color=SAUDI)
    textbox(s, '~2 mm yr⁻¹',
            Inches(5.0), Inches(3.0), Inches(2), Inches(0.5),
            size=20, bold=True, color=SAUDI)
    textbox(s, 'Italy',
            Inches(2.5), Inches(3.6), Inches(2.5), Inches(0.5),
            size=14, bold=True, color=ITALY)
    textbox(s, '~270 mm yr⁻¹',
            Inches(5.0), Inches(3.6), Inches(2.5), Inches(0.5),
            size=20, bold=True, color=ITALY)
    textbox(s, 'Bangladesh',
            Inches(2.5), Inches(4.2), Inches(2.5), Inches(0.5),
            size=14, bold=True, color=BANGLA)
    textbox(s, '~500 mm yr⁻¹',
            Inches(5.0), Inches(4.2), Inches(2.5), Inches(0.5),
            size=20, bold=True, color=BANGLA)
    # 250x callout
    filled_rect(s, Inches(0.5), Inches(5.4), Inches(12.3),
                Inches(0.95), RGBColor(0xFC, 0xEB, 0xE7))
    textbox(s, '→  ~250× increase across the aridity gradient',
            Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.55),
            size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    textbox(s, 'Same feature.   Same model.   Two orders of magnitude in control.',
            Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4),
            size=14, color=GREY, align=PP_ALIGN.CENTER)
    # Take-away bullet at bottom
    textbox(s,
            'Physically: water-limited soils consume rainfall before runoff forms; '
            'humid soils are near saturation and rainfall converts directly.',
            Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5),
            size=13, color=DARK, align=PP_ALIGN.CENTER)
    set_notes(s,
              "This is our key finding. Precipitation is the top driver in all three "
              "countries — but its absolute magnitude differs by two orders of magnitude. "
              "In Saudi Arabia, the mean absolute SHAP of precipitation is just 2 millimetres "
              "per year. In Italy, 270. In Bangladesh, 500. That's a 250-fold difference. "
              "Same feature, same model, same dataset — yet the marginal effect of each "
              "millimetre of rainfall on runoff is essentially zero in Saudi Arabia and "
              "approximately one in Bangladesh. Physically this corresponds to the "
              "water-limited to energy-limited transition.")

    # ---------- 17. SHAP — beeswarm ----------
    s = add_blank(prs)
    apply_chrome(s, 'Dominant runoff drivers (SHAP) — feature distributions', 17)
    add_image(s, 'fig11b_annual_shap_summary.png',
              Inches(0.3), Inches(1.1), width=Inches(12.7))
    bullets = [
        (BANGLA, 'Bangladesh: ΔS ranks #2',
         'saturation-excess regime — antecedent wetness modulates P → R efficiency'),
        (ITALY, 'Italy: ΔS also important (#3)',
         'soil moisture memory matters in transitional climates too'),
    ]
    for i, (col, head, body) in enumerate(bullets):
        y = Inches(5.5 + i * 0.7)
        textbox(s, head, Inches(0.4), y, Inches(5), Inches(0.4),
                size=14, bold=True, color=col)
        textbox(s, body, Inches(5.5), y, Inches(7.5), Inches(0.4),
                size=13, color=DARK)
    set_notes(s,
              "Looking at what ranks second tells us about the mechanism. In Bangladesh, "
              "soil water change is the second most important feature — this is the classic "
              "saturation-excess regime, where antecedent soil wetness modulates how "
              "efficiently rainfall converts to runoff. Italy also has soil water change in "
              "the top three, suggesting soil moisture memory matters in transitional "
              "climates. For Saudi Arabia, all features have very low SHAP magnitude — "
              "consistent with the low predictability we already discussed.")

    # ---------- 18. SHAP — dependence ----------
    s = add_blank(prs)
    apply_chrome(s, 'Dominant runoff drivers (SHAP) — precipitation dependence', 18)
    add_image(s, 'fig13b_annual_shap_dependence_P.png',
              Inches(0.3), Inches(1.1), width=Inches(12.7))
    bullets = [
        (ITALY, 'Italy: threshold response',
         'below ~500 mm yr⁻¹ P barely contributes; above it, SHAP rises steeply'),
        (ACCENT, 'Warmer years suppress SHAP (red points)',
         'higher ET demand reduces P → R efficiency'),
    ]
    for i, (col, head, body) in enumerate(bullets):
        y = Inches(5.5 + i * 0.7)
        textbox(s, head, Inches(0.4), y, Inches(5.2), Inches(0.4),
                size=14, bold=True, color=col)
        textbox(s, body, Inches(5.7), y, Inches(7.3), Inches(0.4),
                size=13, color=DARK)
    set_notes(s,
              "The dependence plot reveals the non-linear shape of precipitation's effect. "
              "Italy shows a clear threshold response — below about 500 millimetres per year, "
              "additional rainfall barely contributes to runoff; above this threshold, the "
              "SHAP contribution rises steeply. The colour shows annual temperature: warmer "
              "years consistently produce lower SHAP at the same precipitation level, "
              "indicating that higher atmospheric demand suppresses the precipitation-to-"
              "runoff conversion efficiency.")

    # ---------- 19. Conclusion divider ----------
    s = add_blank(prs)
    section_divider_chrome(s, 'Conclusion')
    set_notes(s, "To wrap up.")

    # ---------- 20. Remarks ----------
    s = add_blank(prs)
    apply_chrome(s, 'Remarks', 20)
    boxes = [
        ('Runoff partitioning differs systematically across regimes',
         'Bangladesh: highest RC (0.56)\nSaudi: lowest (0.03)\n≈ 20× range'),
        ('Hydroclimatic trends are region-specific',
         'All warm.\nBangladesh dries.\nItaly intensifies ET.'),
        ('Dominant drivers shift across climates',
         'Episodic rainfall (arid)\n→ mixed controls (transitional)\n→ P dominance (humid)\nP SHAP ↑ ~250×'),
        ('SHAP quantifies model behaviour',
         'Not strict causation —\ninterpretations aligned with\nthe Budyko framework for\nphysical credibility.'),
    ]
    box_w = Inches(2.95)
    box_h = Inches(4.5)
    box_top = Inches(1.6)
    accent_cols = [BANGLA, ITALY, ACCENT, GREY]
    for i, (head, body) in enumerate(boxes):
        left = Inches(0.4 + i * 3.15)
        # Background card
        filled_rect(s, left, box_top, box_w, box_h,
                    RGBColor(0xF8, 0xFB, 0xF8))
        # Top accent strip
        filled_rect(s, left, box_top, box_w, Inches(0.18),
                    accent_cols[i])
        # Header
        textbox(s, head, left + Inches(0.15),
                box_top + Inches(0.4), box_w - Inches(0.3),
                Inches(1.5),
                size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        # Body
        textbox(s, body, left + Inches(0.15),
                box_top + Inches(2.0), box_w - Inches(0.3),
                Inches(2.4),
                size=12, color=DARK, align=PP_ALIGN.CENTER)
    set_notes(s,
              "Four take-aways. First, runoff partitioning differs systematically across "
              "regimes. Second, hydroclimatic trends are region-specific — all warm, but "
              "only Bangladesh dries and only Italy shows ET intensification. Third, "
              "dominant drivers shift from episodic rainfall in arid zones to clear "
              "precipitation dominance in humid zones — and the precipitation SHAP magnitude "
              "itself increases by roughly 250 times across the gradient. Fourth, SHAP "
              "quantifies model behaviour rather than strict causation, but our "
              "interpretations are aligned with the Budyko framework for physical credibility.")

    # ---------- 21. Thank you ----------
    s = add_blank(prs)
    # Big triangle accent
    sh = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, Inches(7.5), Inches(0.6),
        Inches(5.83), Inches(6.5),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = DIVIDER_GREEN
    sh.line.fill.background()
    sh.rotation = 90
    filled_rect(s, 0, 0, SLIDE_W, Inches(0.6), BANNER_GREEN)
    filled_rect(s, 0, SLIDE_H - Inches(0.4), SLIDE_W,
                Inches(0.4), BANNER_GREEN)
    textbox(s, 'Thank you!',
            Inches(0.6), Inches(2.8), Inches(7), Inches(1.5),
            size=70, bold=True, color=DARK, font='Cambria')
    textbox(s, 'Questions?',
            Inches(0.6), Inches(4.4), Inches(7), Inches(0.7),
            size=28, color=GREY, font='Cambria')
    set_notes(s,
              "Thank you for listening. We're happy to take questions.")

    # ---------- save ----------
    os.makedirs(os.path.dirname(OUT_PPTX), exist_ok=True)
    prs.save(OUT_PPTX)
    print(f'Saved: {OUT_PPTX}  ({len(prs.slides)} slides)')


if __name__ == '__main__':
    build()
